import os
import time
import hashlib
import pandas as pd
import re
import pdfplumber
from pptx import Presentation
from elasticsearch import Elasticsearch, helpers
from dotenv import load_dotenv
from tqdm import tqdm

# Hàm làm sạch nội dung
def clean_content(text):
    # Loại bỏ dấu câu và chuyển thành chữ thường
    cleaned_text = re.sub(r'\s+', ' ', text)  # Loại bỏ khoảng trắng thừa
    cleaned_text = re.sub(r'[^\w\s]', '', cleaned_text)  # Loại bỏ dấu câu
    return cleaned_text.lower()

# Hàm tạo index với Synonym, autocomplete, và ngram
def create_index_if_not_exists(index_name, es):
    if not es.indices.exists(index=index_name):
        es.indices.create(index=index_name, body={
            "settings": {
                "analysis": {
                    "tokenizer": {
                        "edge_ngram_tokenizer": {
                            "type": "edge_ngram",
                            "min_gram": 1,
                            "max_gram": 25,
                            "token_chars": ["letter", "digit"]
                        }
                    },
                    "analyzer": {
                        "edge_ngram_analyzer": {
                            "type": "custom",
                            "tokenizer": "edge_ngram_tokenizer"
                        }
                    },
                    "filter": {
                        "synonym_filter": {
                            "type": "synonym",
                            "synonyms": [
                                "car, automobile, vehicle",
                                "quick, fast, speedy"
                            ]
                        }
                    }
                }
            },
            "mappings": {
                "properties": {
                    "content": {"type": "text", "analyzer": "edge_ngram_analyzer"},
                    "content_clean": {"type": "text", "analyzer": "standard"},
                    "filename": {"type": "keyword"},
                    "folder": {"type": "keyword"},
                    "content_hash": {"type": "keyword"},
                    "length": {"type": "integer"},
                    "file_extension": {"type": "keyword"},
                    "last_modified": {"type": "date", "format": "epoch_second"},
                    "autocomplete": {
                        "type": "text",
                        "analyzer": "edge_ngram_analyzer"
                    }
                }
            }
        })

def index_documents(base_folder="txt"):
    # Tải biến môi trường
    load_dotenv()

    # Kết nối Elasticsearch
    ES_HOST = os.getenv("ES_HOST", "http://localhost:9200")
    es = Elasticsearch(ES_HOST)

    def hash_content(content):
        """Tạo hash SHA256 cho nội dung file để kiểm tra thay đổi."""
        return hashlib.sha256(content.encode("utf-8")).hexdigest()

    start_time = time.time()

    for root, _, files in os.walk(base_folder):
        folder_name = os.path.basename(root)
        index_name = f"documents_{folder_name.lower()}"
        create_index_if_not_exists(index_name, es)

        batch_documents = []

        for file_name in tqdm(files, desc=f"📂 Đang xử lý ({folder_name})"):
            file_path = os.path.join(root, file_name)
            content = ""

            try:
                if file_name.endswith(".txt"):
                    with open(file_path, "r", encoding="utf-8") as f:
                        content = f.read()

                elif file_name.endswith(".xlsx"):
                    df = pd.read_excel(file_path, sheet_name=None, engine="openpyxl")
                    content = "\n".join([df[sheet].to_string(index=False) for sheet in df])

                elif file_name.endswith(".xls"):
                    try:
                        df = pd.read_excel(file_path, sheet_name=None, engine="xlrd")
                        content = "\n".join([df[sheet].to_string(index=False) for sheet in df])
                    except Exception as xlrd_error:
                        print(f"⚠️ Lỗi xlrd với file {file_name}, thử openpyxl: {xlrd_error}")
                        try:
                            df = pd.read_excel(file_path, sheet_name=None, engine="openpyxl")
                            content = "\n".join([df[sheet].to_string(index=False) for sheet in df])
                        except Exception as openpyxl_error:
                            print(f"❌ Không thể đọc file {file_name}: {openpyxl_error}")
                            continue

                elif file_name.endswith(".pdf"):
                    with pdfplumber.open(file_path) as pdf:
                        content = "\n".join([page.extract_text() for page in pdf.pages])

                elif file_name.endswith(".pptx"):
                    prs = Presentation(file_path)
                    content = "\n".join([slide.shapes[0].text for slide in prs.slides if hasattr(slide.shapes[0], "text")])

            except Exception as e:
                print(f"⚠️ Lỗi đọc file {file_name}: {e}")
                continue

            if not content:
                continue

            content_hash = hash_content(content)
            doc_id = hashlib.sha1(content.encode("utf-8")).hexdigest()

            # Kiểm tra file đã tồn tại chưa (tránh index lại)
            try:
                existing = es.search(index=index_name, body={
                    "query": {"term": {"content_hash.keyword": content_hash}},
                    "size": 1
                })
                if existing["hits"]["hits"]:
                    continue
            except Exception as e:
                print(f"⚠️ Lỗi khi kiểm tra content_hash: {e}")

            batch_documents.append({
                "_op_type": "index",
                "_index": index_name,
                "_id": doc_id,
                "_source": {
                    "content": content,
                    "content_clean": clean_content(content),
                    "filename": file_name,
                    "folder": folder_name,
                    "content_hash": content_hash,
                    "length": len(content),
                    "file_extension": os.path.splitext(file_name)[-1].lower(),
                    "last_modified": int(os.path.getmtime(file_path)),
                    "autocomplete": content[:30]  # Lấy phần đầu của content cho autocomplete
                }
            })

        if batch_documents:
            helpers.bulk(es, batch_documents)
            print(f"✅ Đã index {len(batch_documents)} tài liệu mới trong {folder_name}.")

    end_time = time.time()
    print(f"🚀 Hoàn tất nạp dữ liệu! ⏳ Tổng thời gian: {end_time - start_time:.2f} giây")
    return {"success": 1, "message": "Đánh index thành công!"}

# 👉 Gọi hàm nếu muốn chạy trực tiếp
if __name__ == "__main__":
    index_documents()
